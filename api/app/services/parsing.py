from pydantic import BaseModel
from pathlib import Path
from typing import Callable
import pymupdf
from docx import Document as DocxDocument
from pptx import Presentation

class ParsedDocument(BaseModel):
    text: str
    page_count: int
    parser: str

def clean_extracted_text(text: str) -> str:
    lines = [line.strip() for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n")]
    non_empty_lines = [line for line in lines if line]
    return "\n".join(non_empty_lines)

def extract_pdf_text(file_path: Path) -> ParsedDocument:
    try:
        doc = pymupdf.open(file_path)
    except Exception as exc:
        raise ValueError(f"Failed to open file: {exc}") from exc

    try:
        page_texts: list[str] = []
        for page_index in range(doc.page_count):
            page = doc.load_page(page_index)
            raw_text = page.get_text("text")
            cleaned_text = clean_extracted_text(raw_text)

            if cleaned_text:
                page_texts.append(f"[Page {page_index + 1}]\n{cleaned_text}")
        full_text = "\n\n".join(page_texts).strip()
        page_count = doc.page_count
    finally:
        doc.close()

    if page_count > 0 and len(full_text) < 50:
        raise ValueError("PDF appears scanned or text is not extractable; OCR fallback not implemented yet")

    return ParsedDocument(text=full_text, page_count=page_count, parser="pymupdf")

def extract_docx_text(file_path: Path) -> ParsedDocument:
    try:
        doc = DocxDocument(str(file_path))
    except Exception as exc:
        raise ValueError(f"Failed to open file: {exc}") from exc

    parts: list[str] = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            if row_text.strip():
                parts.append(row_text)

    full_text = clean_extracted_text("\n".join(parts))

    if len(full_text) < 50:
        raise ValueError("Document appears to contain no extractable text")

    return ParsedDocument(text=full_text, page_count=1, parser="python-docx")

def extract_pptx_text(file_path: Path) -> ParsedDocument:
    try:
        presentation = Presentation(str(file_path))
    except Exception as exc:
        raise ValueError(f"Failed to open file: {exc}") from exc

    slide_texts: list[str] = []
    for slide_index, slide in enumerate(presentation.slides):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    run_text = "".join(run.text for run in paragraph.runs)
                    if run_text.strip():
                        parts.append(run_text)
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        parts.append(row_text)
        cleaned_slide = clean_extracted_text("\n".join(parts))
        if cleaned_slide:
            slide_texts.append(f"[Slide {slide_index + 1}]\n{cleaned_slide}")

    full_text = "\n\n".join(slide_texts).strip()
    page_count = len(presentation.slides)

    if page_count > 0 and len(full_text) < 50:
        raise ValueError("Presentation appears to contain no extractable text")

    return ParsedDocument(text=full_text, page_count=page_count, parser="python-pptx")

EXTRACTORS: dict[str, Callable[[Path], ParsedDocument]] = {
    ".pdf": extract_pdf_text,
    ".docx": extract_docx_text,
    ".pptx": extract_pptx_text,
}

def extract_document_text(file_path: Path) -> ParsedDocument:
    extractor = EXTRACTORS.get(file_path.suffix.lower())
    if extractor is None:
        raise ValueError(f"Unsupported file type: {file_path.suffix or 'unknown'}")
    return extractor(file_path)
